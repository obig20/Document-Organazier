"""
Advanced Document Processing Service
Handles document conversion, enhanced OCR, and metadata extraction
"""
import os
import logging
import tempfile
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Union, BinaryIO
from enum import Enum

import fitz  # PyMuPDF
from PIL import Image
import pytesseract
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
import magic

from app.core.config import settings
from app.models.document import Document, DocumentMetadata

logger = logging.getLogger(__name__)

class DocumentType(str, Enum):
    PDF = "application/pdf"
    DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    IMAGE = "image"
    TEXT = "text/plain"
    UNKNOWN = "application/octet-stream"

class DocumentProcessor:
    """Advanced document processing with format conversion and enhanced OCR"""
    
    def __init__(self):
        self.supported_formats = {
            '.pdf': DocumentType.PDF,
            '.docx': DocumentType.DOCX,
            '.xlsx': DocumentType.XLSX,
            '.pptx': DocumentType.PPTX,
            '.txt': DocumentType.TEXT,
            '.jpg': DocumentType.IMAGE,
            '.jpeg': DocumentType.IMAGE,
            '.png': DocumentType.IMAGE,
            '.tiff': DocumentType.IMAGE,
            '.bmp': DocumentType.IMAGE,
        }
    
    def detect_document_type(self, file_path: Union[str, Path, BinaryIO]) -> Tuple[DocumentType, str]:
        """Detect document type using both file extension and content analysis"""
        # First try file extension
        if isinstance(file_path, (str, Path)):
            file_path = Path(file_path)
            ext = file_path.suffix.lower()
            if ext in self.supported_formats:
                return self.supported_formats[ext], ext
        
        # Fall back to content-based detection
        mime = magic.Magic(mime=True)
        if hasattr(file_path, 'read'):
            file_path.seek(0)
            content = file_path.read(1024)
            mime_type = mime.from_buffer(content)
        else:
            mime_type = mime.from_file(str(file_path))
        
        # Map MIME type to our document types
        for ext, doc_type in self.supported_formats.items():
            if doc_type.value == mime_type:
                return doc_type, ext
        
        return DocumentType.UNKNOWN, ''
    
    def extract_text(self, file_path: Union[str, Path], language: str = 'eng+amh+orm') -> str:
        """Extract text from document with enhanced OCR support"""
        doc_type, _ = self.detect_document_type(file_path)
        
        if doc_type == DocumentType.PDF:
            return self._extract_text_from_pdf(file_path, language)
        elif doc_type == DocumentType.DOCX:
            return self._extract_text_from_docx(file_path)
        elif doc_type == DocumentType.XLSX:
            return self._extract_text_from_xlsx(file_path)
        elif doc_type == DocumentType.PPTX:
            return self._extract_text_from_pptx(file_path)
        elif doc_type == DocumentType.IMAGE:
            return self._extract_text_from_image(file_path, language)
        elif doc_type == DocumentType.TEXT:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        else:
            raise ValueError(f"Unsupported document type: {doc_type}")
    
    def _extract_text_from_pdf(self, file_path: Union[str, Path], language: str) -> str:
        """Extract text from PDF with OCR fallback"""
        text = ""
        try:
            # First try to extract text directly
            with fitz.open(file_path) as doc:
                for page in doc:
                    text += page.get_text()
            
            # If no text found, try OCR
            if not text.strip():
                text = self._ocr_pdf(file_path, language)
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {e}")
            text = self._ocr_pdf(file_path, language)
        
        return text
    
    def _extract_text_from_docx(self, file_path: Union[str, Path]) -> str:
        """Extract text from DOCX file"""
        doc = DocxDocument(file_path)
        return "\n\n".join(paragraph.text for paragraph in doc.paragraphs)
    
    def _extract_text_from_xlsx(self, file_path: Union[str, Path]) -> str:
        """Extract text from XLSX file"""
        wb = load_workbook(file_path, read_only=True, data_only=True)
        text_parts = []
        
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    text_parts.append(row_text)
        
        return "\n".join(text_parts)
    
    def _extract_text_from_pptx(self, file_path: Union[str, Path]) -> str:
        """Extract text from PPTX file"""
        prs = Presentation(file_path)
        text_parts = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
        
        return "\n\n".join(text_parts)
    
    def _extract_text_from_image(self, file_path: Union[str, Path], language: str) -> str:
        """Extract text from image using Tesseract OCR"""
        try:
            return pytesseract.image_to_string(
                str(file_path),
                lang=language,
                config='--psm 6 --oem 3'  # Assume a uniform block of text
            )
        except Exception as e:
            logger.error(f"Error performing OCR on image: {e}")
            return ""
    
    def _ocr_pdf(self, file_path: Union[str, Path], language: str) -> str:
        """Perform OCR on PDF pages"""
        text_parts = []
        
        try:
            # Convert PDF to images and perform OCR on each page
            with fitz.open(file_path) as doc:
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    pix = page.get_pixmap()
                    
                    # Convert to PIL Image
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    
                    # Perform OCR
                    page_text = pytesseract.image_to_string(
                        img,
                        lang=language,
                        config='--psm 6 --oem 3'  # Assume a uniform block of text
                    )
                    
                    if page_text.strip():
                        text_parts.append(page_text)
        except Exception as e:
            logger.error(f"Error during PDF OCR: {e}")
        
        return "\n\n".join(text_parts)
    
    def convert_document(
        self,
        source_path: Union[str, Path],
        target_format: str,
        output_dir: Optional[Union[str, Path]] = None
    ) -> Path:
        """Convert document to target format"""
        # Implementation for document conversion
        # This is a placeholder - actual implementation would depend on the conversion requirements
        raise NotImplementedError("Document conversion will be implemented in a future update")
    
    def extract_metadata(self, file_path: Union[str, Path]) -> DocumentMetadata:
        """Extract metadata from document"""
        path = Path(file_path)
        stat = path.stat()
        
        return DocumentMetadata(
            file_size=stat.st_size,
            file_type=path.suffix.lower().lstrip('.'),
            created_date=stat.st_ctime,
            modified_date=stat.st_mtime,
            mime_type=magic.from_file(str(path), mime=True)
        )
    
    def process_document(self, file_path: Union[str, Path], language: str = 'eng+amh+orm') -> Document:
        """Process document and return Document object"""
        path = Path(file_path)
        
        # Extract text and metadata
        text = self.extract_text(path, language)
        metadata = self.extract_metadata(path)
        
        # Create Document object
        return Document(
            filename=path.name,
            original_path=str(path.absolute()),
            content=text,
            metadata=metadata,
            is_processed=True,
            processing_status="completed"
        )

# Singleton instance
document_processor = DocumentProcessor()
